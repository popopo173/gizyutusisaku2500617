from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
from io import BytesIO
from flask_cors import CORS
from PIL import Image

app = Flask(__name__)
CORS(app)

@app.route("/export", methods=["POST"])
def export_pptx():
    """
    フロントエンドから送られてきた画像URL3つを背景画像として、
    PowerPointファイルを生成し、ダウンロード用に送信する。
    スライド構成:
        - スライド1: タイトルスライド（中央にタイトル）
        - スライド2: 本文スライド（タイトルと本文を左寄せ表示）
        - スライド3: 章区切りスライド（セクションタイトルとサブタイトル）
    """
    
    data = request.json
    images = data.get("images", [])

    if len(images) != 3:
        return {"error": "3 images are required"}, 400

    prs = Presentation()

    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # --- スライド1: タイトルスライド（中央にテキスト） ---
    slide_layout = prs.slide_layouts[6]  # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(prs, slide1, images[0])
    add_textbox(slide1, "サンプルタイトル", Inches(3), Inches(2), Inches(5), Inches(0.7), font_size=40)

    # --- スライド2: 本文スライド（タイトル＋本文） ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blankに変更
    set_background(prs, slide2, images[1])
    add_textbox(slide2, "本文スライドのタイトル", Inches(0.5), Inches(0), Inches(8), Inches(1), font_size=28)
    add_textbox(slide2, "ここに本文内容を入力", Inches(1), Inches(1), Inches(8), Inches(4), font_size=20)

    # --- スライド3: 章区切りスライド（セクションタイトル＋サブタイトル） ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blankに変更
    set_background(prs, slide3, images[2])
    add_textbox(slide3, "章区切りスライド", Inches(3), Inches(1.5), Inches(5), Inches(0.7), font_size=32)
    add_textbox(slide3, "セクションの紹介文", Inches(3), Inches(3), Inches(5), Inches(0.7), font_size=24)

    # --- 書き出し ---
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, download_name="slides.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

def set_background(prs, slide, image_url):
    """
    スライドの背景として画像を挿入する。
    背景画像を16:9にトリミングして設定
    """
    response = requests.get(image_url)
    original_img = Image.open(BytesIO(response.content))

    width, height = original_img.size
    target_ratio = 16 / 9
    new_height = int(width / target_ratio)

    if new_height > height:
        # 横長すぎる場合
        new_width = int(height * target_ratio)
        left = (width - new_width) // 2
        top = 0
        right = left + new_width
        bottom = height
    else:
        # 縦長すぎる場合
        top = (height - new_height) // 2
        left = 0
        bottom = top + new_height
        right = width

    cropped = original_img.crop((left, top, right, bottom))
    cropped_stream = BytesIO()
    cropped.save(cropped_stream, format='PNG')
    cropped_stream.seek(0)

    slide.shapes.add_picture(
        cropped_stream,
        0, 0,
        width=prs.slide_width,
        height=prs.slide_height
    )

def add_textbox(slide, text, left, top, width, height, font_size=24):
    """
    指定された位置とサイズにテキストボックスを配置し、文字列を表示する。
    """
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)

if __name__ == "__main__":
    app.run(port=5000)
