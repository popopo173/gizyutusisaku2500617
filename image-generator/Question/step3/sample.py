from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
from io import BytesIO
from flask_cors import CORS
from PIL import Image

app = Flask(__name__)
CORS(app)

@app.route("/export", methods=["POST"])
def export_sample():
    """
    発展課題用サンプル:
        - スライドに図形（四角形と円）を追加
        - 小さい画像を任意の位置に配置
        - 箇条書きテキストを作成
    """

    data = request.json
    images = data.get("images", [])

    if len(images) < 1:
        return {"error": "少なくとも1枚の画像が必要です"}, 400

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # --- スライド1 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 1. 図形を追加（四角形 & 円）
    from pptx.enum.shapes import MSO_SHAPE
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
    )
    slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(2)
    )

    # 2. 任意の位置に小さい画像を配置
    response = requests.get(images[0])
    icon_stream = BytesIO(response.content)
    slide.shapes.add_picture(icon_stream, Inches(7), Inches(1), Inches(2), Inches(2))

    # 3. 箇条書きテキストを作成
    textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(2))
    text_frame = textbox.text_frame
    text_frame.clear()

    p1 = text_frame.add_paragraph()
    p1.text = "これは1つ目の箇条書き項目です"
    p1.level = 0  # レベル0 → 第一階層

    p2 = text_frame.add_paragraph()
    p2.text = "これは2つ目の箇条書き項目です"
    p2.level = 0

    p3 = text_frame.add_paragraph()
    p3.text = "さらに詳細なサブ項目"
    p3.level = 1  # レベル1 → インデント付き

    # --- 書き出し ---
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(
        output,
        as_attachment=True,
        download_name="sample_tasks.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == "__main__":
    app.run(port=5000)
